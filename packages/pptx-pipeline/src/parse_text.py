"""Extract structured text content from PPTX for AI processing."""

from pathlib import Path

from pptx import Presentation


def extract_text_content(pptx_path: Path) -> list[dict[str, str]]:
    """Extract text from each slide as structured data.

    Returns a list of dicts with 'slide_number', 'title', 'body', and 'notes' keys.
    Useful for feeding into AI quiz generation.
    """
    prs = Presentation(str(pptx_path))
    result = []

    for i, slide in enumerate(prs.slides):
        slide_text: dict[str, str] = {
            "slide_number": str(i + 1),
            "title": "",
            "body": "",
            "notes": "",
        }

        body_parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text.strip()
            if not text:
                continue

            # First text shape with content is likely the title
            if not slide_text["title"]:
                slide_text["title"] = text
            else:
                body_parts.append(text)

        slide_text["body"] = "\n".join(body_parts)

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_text["notes"] = slide.notes_slide.notes_text_frame.text

        result.append(slide_text)

    return result
