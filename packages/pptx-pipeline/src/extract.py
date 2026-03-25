"""Extract slide data from PPTX using python-pptx."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def extract_slide_data(pptx_path: Path) -> list[dict]:
    """Extract text, notes, and shape data from each slide."""
    prs = Presentation(str(pptx_path))
    slides_data = []

    for slide_index, slide in enumerate(prs.slides):
        slide_info: dict = {
            "index": slide_index,
            "title": "",
            "text": "",
            "notes": "",
            "shapes": [],
        }

        texts = []
        for shape in slide.shapes:
            # Extract shape geometry
            shape_data = {
                "type": _get_shape_type(shape),
                "x": _emu_to_px(shape.left) if shape.left else 0,
                "y": _emu_to_px(shape.top) if shape.top else 0,
                "width": _emu_to_px(shape.width) if shape.width else 0,
                "height": _emu_to_px(shape.height) if shape.height else 0,
            }

            if shape.has_text_frame:
                text = shape.text_frame.text
                texts.append(text)
                shape_data["content"] = text

                # Check for title
                if shape.shape_id == 1 or (hasattr(shape, "placeholder_format") and shape.placeholder_format is not None):
                    if not slide_info["title"]:
                        slide_info["title"] = text

            slide_info["shapes"].append(shape_data)

        slide_info["text"] = "\n".join(texts)

        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_info["notes"] = slide.notes_slide.notes_text_frame.text

        slides_data.append(slide_info)

    return slides_data


def _get_shape_type(shape) -> str:
    """Determine shape type string."""
    if shape.shape_type is not None:
        name = str(shape.shape_type).lower()
        if "picture" in name or "image" in name:
            return "image"
        if "text" in name:
            return "text"
    if shape.has_text_frame:
        return "text"
    return "shape"


def _emu_to_px(emu_value: int, dpi: int = 96) -> int:
    """Convert EMU (English Metric Units) to pixels."""
    return round(emu_value / 914400 * dpi)
